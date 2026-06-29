"""
PDF アプリ 仕様概要書 (非エンジニア向け) PowerPoint 生成スクリプト
出力: docs/app_overview.pptx
"""

from pptx import Presentation
from pptx.util import Emu, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import os

# ── 定数 ─────────────────────────────────────────────────────────────────────
RED    = RGBColor(0xCC, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
GRAY_L = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_D = RGBColor(0x44, 0x44, 0x44)

SLIDE_W = 12192000   # EMU  (16:9)
SLIDE_H =  6858000   # EMU

FONT    = "Meiryo"
HEADER_H = Cm(1.2)

OUT_PATH = os.path.join(os.path.dirname(__file__), "app_overview.pptx")

# ── ヘルパー ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()   # no line by default
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_name=FONT, font_size=Pt(14), bold=False,
                color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_header_bar(slide, title, page_no=None):
    """スライド上部の赤ヘッダーバーを追加（タイトルスライド以外）"""
    bar = add_rect(slide, 0, 0, SLIDE_W, HEADER_H, fill_rgb=RED)
    tb = slide.shapes.add_textbox(Cm(0.4), 0, SLIDE_W - Cm(1), HEADER_H)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE
    # ページ番号
    if page_no is not None:
        pn = add_textbox(
            slide, str(page_no),
            SLIDE_W - Cm(2), SLIDE_H - Cm(0.8),
            Cm(1.5), Cm(0.7),
            font_size=Pt(10), color=GRAY_D,
            align=PP_ALIGN.RIGHT
        )


def add_table(slide, headers, rows, x, y, w, h, col_widths=None):
    """表を追加。ヘッダー行: 赤背景・白文字、交互行: 薄グレー"""
    rows_total = 1 + len(rows)
    cols_total = len(headers)
    table = slide.shapes.add_table(rows_total, cols_total, x, y, w, h).table

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw

    def set_cell(cell, text, bg=None, fg=BLACK, bold=False, font_size=Pt(12)):
        cell.text = text
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = fg
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        else:
            cell.fill.background()

    # ヘッダー行
    for ci, h_text in enumerate(headers):
        set_cell(table.cell(0, ci), h_text, bg=RED, fg=WHITE, bold=True)

    # データ行
    for ri, row in enumerate(rows):
        bg = GRAY_L if ri % 2 == 0 else WHITE
        for ci, cell_text in enumerate(row):
            set_cell(table.cell(ri + 1, ci), cell_text, bg=bg)

    return table


def set_slide_bg(slide, color_rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb


# ── スライド生成 ──────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """タイトルスライド"""
    sl = blank_slide(prs)
    set_slide_bg(sl, RED)
    # メインタイトル
    add_textbox(sl, "PDF アプリ 仕様概要書",
                Cm(1.5), Cm(2.2), Cm(28), Cm(2.5),
                font_size=Pt(40), bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    # サブタイトル
    add_textbox(sl, "〜 システムを熟知していない方向け 〜",
                Cm(1.5), Cm(4.8), Cm(28), Cm(1.2),
                font_size=Pt(22), bold=False, color=WHITE,
                align=PP_ALIGN.CENTER)
    # 日付
    add_textbox(sl, "2026-06-23",
                Cm(1.5), Cm(6.3), Cm(28), Cm(0.8),
                font_size=Pt(14), color=WHITE,
                align=PP_ALIGN.CENTER)


def slide_02_toc(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "目次", 2)
    items = [
        "1.  アプリ概要 .............................................................. 3",
        "2.  UIの自由度（固定 vs 変更可能） ................................ 5",
        "3.  操作・インタラクション .............................................. 7",
        "4.  PDFビューア 機能詳細 ................................................ 9",
        "5.  補足（テーマ／言語切替・テスト用機能） ................. 11",
    ]
    y = HEADER_H + Cm(0.8)
    for item in items:
        add_textbox(sl, item,
                    Cm(2), y, Cm(27), Cm(0.9),
                    font_size=Pt(16), color=BLACK)
        y += Cm(0.85)


def slide_03_overview(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "1. アプリ概要", 3)

    # 説明文
    desc = (
        "本アプリは、PDFコンテンツをダウンロード・閲覧するためのAndroid向けアプリです。\n"
        "コンテンツ一覧から目的のPDFをダウンロードし、高機能なPDFビューアで閲覧できます。\n"
        "テーマ・言語の切替にも対応しており、ユーザーの好みに合わせた表示が可能です。"
    )
    add_textbox(sl, desc,
                Cm(1), HEADER_H + Cm(0.5), Cm(29), Cm(1.8),
                font_size=Pt(13), color=BLACK)

    # 基本スペック表
    headers = ["項目", "内容"]
    rows = [
        ["対応OS",       "Android"],
        ["通信",         "必要（初回ダウンロード時のみ。閲覧はオフライン可）"],
        ["ストレージ",   "端末内 Documents フォルダに保存"],
        ["表示言語",     "日本語・英語（アプリ内で切替可能）"],
        ["テーマ",       "ライト / ダーク / システム連動（アプリ内で切替可能）"],
    ]
    col_w = [Cm(4.5), Cm(22)]
    add_table(sl, headers, rows,
              Cm(1), HEADER_H + Cm(2.5), Cm(27), Cm(3.2),
              col_widths=col_w)


def slide_04_flow(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "画面フロー", 4)

    # 3つのボックスと矢印
    boxes = [
        (Cm(1.5),  Cm(2.5), Cm(7.5), Cm(3.0), "① コンテンツ一覧",
         "・コンテンツ一覧の確認\n・PDFのダウンロード\n・テーマ・言語切替"),
        (Cm(11.5), Cm(2.5), Cm(7.5), Cm(3.0), "② PDFビューア",
         "・PDFの閲覧\n・ページ操作・ズーム\n・ブックマーク・検索"),
        (Cm(21.5), Cm(2.5), Cm(7.5), Cm(3.0), "③ WebView",
         "・PDF内URLリンクを\n　アプリ内で表示\n・（自動遷移・手動戻れる）"),
    ]
    for bx, by, bw, bh, title, desc in boxes:
        # 外枠
        add_rect(sl, bx, by, bw, bh, fill_rgb=GRAY_L, line_rgb=RED, line_width_pt=1.5)
        # タイトルバー
        add_rect(sl, bx, by, bw, Cm(0.75), fill_rgb=RED)
        add_textbox(sl, title, bx + Cm(0.2), by, bw - Cm(0.4), Cm(0.75),
                    font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_textbox(sl, desc, bx + Cm(0.3), by + Cm(0.85), bw - Cm(0.6), bh - Cm(1.0),
                    font_size=Pt(12), color=BLACK)

    # 矢印テキスト
    for ax in [Cm(9.3), Cm(19.3)]:
        add_textbox(sl, "▶", ax, Cm(3.6), Cm(2.0), Cm(1.0),
                    font_size=Pt(26), color=RED, align=PP_ALIGN.CENTER)

    # 注記
    add_textbox(sl, "※ WebView は PDF 内の URL リンクをタップした場合のみ表示されます。",
                Cm(1), Cm(6.0), Cm(29), Cm(0.6),
                font_size=Pt(11), color=GRAY_D)


def slide_05_ui_fixed(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "2. UIの自由度 ① — 固定部分（変更できないもの）", 5)

    headers = ["項目", "具体例・詳細"]
    rows = [
        ["レイアウト構造",     "画面内の各要素（ヘッダー・コンテンツ・フッター）の配置・余白"],
        ["ナビゲーション方式", "コンテンツ一覧 → PDFビューア → WebView の一方向フロー"],
        ["ジェスチャー操作",   "横スワイプ・ピンチズーム・タップ等の操作方法"],
        ["アイコン配置",       "AppBar 右上のボタン群（テーマ切替・言語切替・表示モード切替）"],
        ["ストレージ保存先",   "端末内 Documents フォルダ（変更不可）"],
        ["ファイル名形式",     "ダウンロードファイルは「ID_言語_タイトル.pdf」形式で保存"],
    ]
    col_w = [Cm(5.5), Cm(21)]
    add_table(sl, headers, rows,
              Cm(1), HEADER_H + Cm(0.5), Cm(27.5), Cm(4.5),
              col_widths=col_w)

    add_textbox(sl, "※ 上記の変更には開発・リリース作業が必要です。",
                Cm(1), SLIDE_H - Cm(1.2), Cm(20), Cm(0.6),
                font_size=Pt(11), color=GRAY_D)


def slide_06_ui_variable(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "2. UIの自由度 ② — 変更可能部分", 6)

    headers = ["変更できるもの", "誰が変えるか", "変え方"]
    rows = [
        ["テーマ（ライト/ダーク）",   "ユーザー",   "アプリ内設定ボタン（即時反映・設定保持）"],
        ["表示言語（日本語/英語）",   "ユーザー",   "アプリ内設定ボタン（即時反映・設定保持）"],
        ["表示モード（リスト/グリッド）", "ユーザー", "AppBar 右上の切替ボタン"],
        ["コンテンツ一覧の内容",      "管理者",     "assets/contents.json を編集（タイトル・説明・URL等）"],
        ["ブランドカラー（赤 #CC0000）", "開発者",  "コード上の seedColor を変更（再ビルド必要）"],
        ["対応言語の追加",            "開発者",     "多言語対応ファイル追加 + 開発作業"],
    ]
    col_w = [Cm(6.5), Cm(5.0), Cm(16.0)]
    add_table(sl, headers, rows,
              Cm(1), HEADER_H + Cm(0.4), Cm(28), Cm(5.0),
              col_widths=col_w)


def slide_07_ux_list(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "3. インタラクション ① — コンテンツ一覧画面", 7)

    headers = ["操作", "結果・動作"]
    rows = [
        ["ダウンロードボタン タップ",       "ダウンロード開始 → 進捗バーを表示 → 完了後「開く」ボタンに変化"],
        ["キャンセルボタン タップ（DL中）",  "ダウンロードを中止（ファイルは保存されない）"],
        ["「開く」ボタン タップ",            "PDFビューア画面へ遷移してPDFを表示"],
        ["削除ボタン タップ",               "確認ダイアログ表示 → OK でファイル削除・未DL状態に戻る"],
        ["リスト/グリッド 切替ボタン",       "リスト表示（テキスト中心）⇔ グリッド表示（サムネイル中心）"],
        ["テーマ切替ボタン",                "ライト / ダーク / システム連動 の選択ダイアログを表示"],
        ["言語切替ボタン",                  "日本語 / 英語 の選択ダイアログを表示"],
    ]
    col_w = [Cm(8.0), Cm(19.5)]
    add_table(sl, headers, rows,
              Cm(1), HEADER_H + Cm(0.4), Cm(28), Cm(5.1),
              col_widths=col_w)


def slide_08_ux_viewer(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "3. インタラクション ② — PDFビューア画面", 8)

    headers = ["操作", "結果・動作"]
    rows = [
        ["横スワイプ",                    "前後のページへ移動"],
        ["ピンチアウト（2本指で広げる）",  "拡大（最大 5 倍）"],
        ["ピンチイン（2本指で縮める）",    "縮小（標準倍率まで）"],
        ["拡大中に横スワイプ",            "ページ送りは無効 → 画面のパン（スクロール）操作になる"],
        ["画面タップ（1回）",             "上部バー・下部サムネイルの 表示 ⇔ 非表示 切替"],
        ["ブックマークボタン タップ",      "現在ページをブックマーク追加 / 解除（アプリ終了後も保持）"],
        ["サムネイル タップ",             "タップしたページへジャンプ"],
        ["PDF内URLリンク タップ",         "アプリ内 WebView でウェブページを表示"],
        ["PDF内ページリンク タップ",       "リンク先のページへジャンプ"],
    ]
    col_w = [Cm(8.5), Cm(19.0)]
    add_table(sl, headers, rows,
              Cm(1), HEADER_H + Cm(0.4), Cm(28), Cm(5.5),
              col_widths=col_w)


def slide_09_features(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "4. PDFビューア 機能一覧", 9)

    features = [
        ("ページナビゲーション", "横スワイプ + 下部サムネイルストリップでページ操作"),
        ("ピンチズーム",         "最大 5 倍まで拡大。拡大中はスワイプがパン操作に切替"),
        ("ブックマーク",         "ページ単位で保存。アプリ終了後も維持"),
        ("キーワード検索",       "全ページを対象にテキスト検索。ヒット箇所をハイライト＆前後ナビ"),
        ("目次（アウトライン）", "PDFに目次が含まれる場合、タップでページジャンプ"),
        ("ダークモード対応",     "PDFページの色を反転表示。目に優しい夜間閲覧が可能"),
        ("PDF内リンク対応",      "URLリンクは内蔵WebViewで、ページリンクはジャンプで処理"),
        ("サムネイル一覧",       "画面下部のストリップで全ページを俯瞰。ブックマークページは強調表示"),
        ("UI自動非表示",         "画面タップでバー類を非表示にし、PDFの閲覧領域を最大化"),
    ]

    # 3列レイアウト（3×3グリッド）
    cols = 3
    card_w = Cm(9.2)
    card_h = Cm(1.45)
    gap_x  = Cm(0.4)
    gap_y  = Cm(0.35)
    start_x = Cm(0.6)
    start_y = HEADER_H + Cm(0.5)

    for i, (name, desc) in enumerate(features):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        # カード背景
        add_rect(sl, x, y, card_w, card_h, fill_rgb=GRAY_L, line_rgb=RED, line_width_pt=0.8)
        # タイトル
        add_textbox(sl, name, x + Cm(0.2), y + Cm(0.05), card_w - Cm(0.4), Cm(0.55),
                    font_size=Pt(12), bold=True, color=RED)
        # 説明
        add_textbox(sl, desc, x + Cm(0.2), y + Cm(0.58), card_w - Cm(0.4), Cm(0.8),
                    font_size=Pt(10), color=BLACK)


def slide_10_drawer(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "4. サイドドロワー 詳細（3タブ構成）", 10)

    # 説明文
    add_textbox(sl,
                "PDFビューア左上のメニューボタンをタップすると、サイドドロワーが開きます。\n"
                "目次・ブックマーク・キーワード検索の3タブを切り替えて使用できます。",
                Cm(1), HEADER_H + Cm(0.4), Cm(29), Cm(1.0),
                font_size=Pt(13), color=BLACK)

    headers = ["タブ名", "アイコン", "機能概要"]
    rows = [
        ["目次\n（アウトライン）",
         "≡ リスト",
         "PDFに含まれる目次（章立て）を一覧表示。\n各項目をタップすると該当ページへジャンプ。\n※ 目次のないPDFでは「目次がありません」と表示。"],
        ["ブックマーク",
         "🔖 ブックマーク",
         "ブックマークしたページの一覧を表示。\nタップで該当ページへジャンプ、削除ボタンで解除。\nブックマークはファイルごとに個別保存される。"],
        ["キーワード検索",
         "🔍 検索",
         "検索ワードを入力してPDF全ページをテキスト検索。\nヒット箇所は画面上でハイライト表示され、\n「前へ／次へ」ボタンで順番にナビゲートできる。"],
    ]
    col_w = [Cm(4.0), Cm(4.5), Cm(19.0)]
    add_table(sl, headers, rows,
              Cm(1), HEADER_H + Cm(1.5), Cm(28), Cm(4.5),
              col_widths=col_w)


def slide_11_settings(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "5. 補足 — テーマ・言語切替", 11)

    add_textbox(sl, "切替方法",
                Cm(1), HEADER_H + Cm(0.5), Cm(20), Cm(0.7),
                font_size=Pt(16), bold=True, color=RED)

    add_textbox(sl,
                "コンテンツ一覧画面の右上にあるボタンからいつでも切り替えられます。",
                Cm(1), HEADER_H + Cm(1.3), Cm(29), Cm(0.7),
                font_size=Pt(13), color=BLACK)

    headers = ["設定項目", "選択肢", "特記事項"]
    rows = [
        ["テーマ", "ライト / ダーク / システム連動",
         "アプリを閉じても設定は保持される。\n「システム連動」は端末のダークモード設定に追従。"],
        ["表示言語", "日本語 / 英語",
         "アプリを閉じても設定は保持される。\nコンテンツ内容（タイトル・説明文）も言語に応じて切替。"],
    ]
    col_w = [Cm(4.5), Cm(7.0), Cm(16.0)]
    add_table(sl, headers, rows,
              Cm(1), HEADER_H + Cm(2.2), Cm(28), Cm(2.8),
              col_widths=col_w)

    add_textbox(sl, "※ ブックマーク・ダウンロード済みファイルはテーマ・言語変更の影響を受けません。",
                Cm(1), SLIDE_H - Cm(1.2), Cm(29), Cm(0.6),
                font_size=Pt(11), color=GRAY_D)


def slide_12_misc(prs):
    sl = blank_slide(prs)
    add_header_bar(sl, "5. 補足 — テスト用機能 / コンテンツ差し替え", 12)

    # テスト用機能
    add_textbox(sl, "🧪 テスト用：ストレージ初期化ボタン",
                Cm(1), HEADER_H + Cm(0.4), Cm(29), Cm(0.7),
                font_size=Pt(15), bold=True, color=RED)

    add_textbox(sl,
                "コンテンツ一覧画面の最下部（フッター）に「ストレージを初期化 (テスト用)」ボタンがあります。\n"
                "タップすると確認ダイアログが表示され、OK を選択するとダウンロード済みPDF全件＋ブックマーク等のデータがすべて削除されます。\n"
                "⚠️ この操作は取り消せません。テスト・検証目的以外での使用は避けてください。",
                Cm(1), HEADER_H + Cm(1.2), Cm(29), Cm(1.8),
                font_size=Pt(12), color=BLACK)

    # 区切り線
    add_rect(sl, Cm(1), HEADER_H + Cm(3.2), Cm(28.5), Cm(0.03), fill_rgb=GRAY_D)

    # コンテンツ差し替え
    add_textbox(sl, "📄 コンテンツ一覧の差し替え方法",
                Cm(1), HEADER_H + Cm(3.4), Cm(29), Cm(0.7),
                font_size=Pt(15), bold=True, color=RED)

    add_textbox(sl,
                "表示されるコンテンツ（タイトル・説明・カテゴリー・ダウンロードURL等）は\n"
                "アプリに同梱された「assets/contents.json」ファイルで管理されています。\n"
                "このファイルを編集してアプリを再ビルドすることで、コンテンツ一覧を自由に差し替えられます。",
                Cm(1), HEADER_H + Cm(4.2), Cm(29), Cm(1.8),
                font_size=Pt(12), color=BLACK)


# ── メイン ────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_toc(prs)
    slide_03_overview(prs)
    slide_04_flow(prs)
    slide_05_ui_fixed(prs)
    slide_06_ui_variable(prs)
    slide_07_ux_list(prs)
    slide_08_ux_viewer(prs)
    slide_09_features(prs)
    slide_10_drawer(prs)
    slide_11_settings(prs)
    slide_12_misc(prs)

    prs.save(OUT_PATH)
    print(f"✅ 保存完了: {OUT_PATH}")


if __name__ == "__main__":
    main()
